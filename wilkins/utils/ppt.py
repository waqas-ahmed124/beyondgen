from pptx import Presentation
from io import BytesIO
from PIL import Image
import logging
from docx.shared import Pt
from pptx.dml.color import RGBColor

logger = logging.getLogger(__name__)


class PowerPointCreator:
    def __init__(self, submissions, template_path):
        self.submissions = submissions
        self.template_path = template_path
        self.placeholder_mappings = {
            21: "Picture",  # Picture placeholder
            22: "market",  # Text placeholder
            23: "location_description",  # Text placeholder
            24: "unit",  # Text placeholder
            25: "media_type",  # Text placeholder for aggregated data
        }

    def add_slide(self, prs, row, image):
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)

        unit = row["unit"]
        # Iterate over all shapes in the slide
        for shape in slide.shapes:
            if not shape.is_placeholder:
                continue

            placeholder_idx = shape.placeholder_format.idx
            placeholder_text = self.placeholder_mappings.get(
                placeholder_idx, "Unknown Placeholder"
            )
            text_frame = shape.text_frame
            text_frame.clear()  # Clear existing text
            newline_check = 0  # A check to add a new line after a few values just to make the slides look a little better
            if placeholder_idx == 25:  # Special handling for the 25th placeholder
                for key in [
                    "media_type",
                    "facing",
                    "size",
                    "is_illuminated",
                    "availability_start",
                    "availability_end",
                    "spot_length_secs",
                    "a18_weekly_impressions",
                    "four_week_media_cost",
                    "installation_cost",
                    "production_cost",
                ]:
                    value = str(row.get(key, "N/A"))
                    if key == "is_illuminated":
                        key = "Illuminated"
                    elif key == "spot_length_secs":
                        key = "Spot Length"
                    elif key == "media_type":
                        key = "Media Type"
                    elif key == "facing":
                        key = "Facing"
                    elif key == "availability_start":
                        key = "Availability Start"
                    elif key == "availability_end":
                        key = "Availability End"
                    elif key == "a18_weekly_impressions":
                        key = "Weekly Impressions"
                    elif key == "four_week_media_cost":
                        key = "4 week Media Cost"
                    elif key == "installation_cost":
                        key = "Installation Cost"
                    elif key == "production_cost":
                        key = "Production Cost"
                    
                    

                    p = (
                        text_frame.add_paragraph()
                    )  # Create a new paragraph for each detail
                    run = p.add_run()
                    run.text = f"{key}: "  # Set the key
                    run.font.bold = True  # Bold for keys

                    run = p.add_run()
                    if newline_check == 2:
                        run.text = f"{value}\n"  # Set the value with a newline
                        newline_check = 0
                    else:
                        run.text = f"{value}"  # Set the value without a newline
                        newline_check += 1
                    run.font.bold = False  # Normal font for values
            else:
                # Set the text for the other placeholders
                value_text = str(row.get(placeholder_text, "N/A"))
                p = text_frame.paragraphs[0]
                # Check if the placeholder text is 'Location Description'
                if placeholder_text == "location_description":
                    run = p.add_run()
                    # Adding 'Location: ' in blue and bold
                    run.text = "Location: "
                    run.font.color.rgb = RGBColor(0, 0, 255)  # Blue color
                    run.font.bold = True

                    def truncate(description):
                        max_length = 100
                        if len(description) <= max_length:
                            return description
                        else:
                            period_index = description.find(".", 0, max_length)
                            if period_index != -1:
                                return description[
                                    : period_index + 1
                                ]  # Includes the period
                            else:
                                comma_index = description.rfind(",", 0, max_length)
                                if comma_index != -1:
                                    return description[
                                        :comma_index
                                    ]  # Excludes the comma
                                else:
                                    # Truncate to max_length if no suitable punctuation found
                                    return description[:max_length]

                    # Adding value_text normally
                    value_run = p.add_run()
                    value_run.text = truncate(value_text)
                    value_run.font.bold = False  # Normal, not bold

                    p.font.size = Pt(14)  # setting font size
                else:
                    run = p.add_run()
                    run.text = f"{placeholder_text}: "
                    run.font.color.rgb = RGBColor(0, 0, 255)  # Blue color
                    run.font.bold = True

                    value_run = p.add_run()
                    value_run.text = value_text
                    value_run.font.bold = True
            # Handle image placeholder separately
            if (
                shape.placeholder_format.idx == 21 and image
            ):  # 21 is the placeholder index for images
                try:
                    if image.mode == "RGBA":
                        image = image.convert("RGB")

                    image_stream = BytesIO()
                    image.save(image_stream, format="JPEG")
                    image_stream.seek(0)

                    # Clear the placeholder content
                    sp = shape._sp
                    sp.getparent().remove(sp)

                    # Add the image at the placeholder's position
                    img_width, img_height = image.size
                    img_aspect_ratio = img_width / img_height
                    placeholder_aspect_ratio = shape.width / shape.height

                    if img_aspect_ratio > placeholder_aspect_ratio:
                        scale_height = int(shape.width / img_aspect_ratio)
                        top = int(shape.top + (shape.height - scale_height) / 2)
                        slide.shapes.add_picture(
                            image_stream, shape.left, top, shape.width, scale_height
                        )
                    else:
                        scale_width = int(shape.height * img_aspect_ratio)
                        left = int(shape.left + (shape.width - scale_width) / 2)
                        slide.shapes.add_picture(
                            image_stream, left, shape.top, scale_width, shape.height
                        )

                except Exception as e:
                    logging.error(f"Error processing image for Unit {unit}: {e}")

    def create_presentation(self, path, azure_storage, container_name="attachments"):
        prs = Presentation(self.template_path)
        for idx, row in self.submissions.iterrows():
            image = row["image_id"]
            if image is not None:
                # Fetch the image from Azure and update the 'image' variable
                image = azure_storage.get_image_as_pil(container_name, image)
            
            self.add_slide(prs, row, image)

        prs.save(path)
        logging.info("PowerPoint created successfully!")
